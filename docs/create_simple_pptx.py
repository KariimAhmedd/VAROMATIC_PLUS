import os
from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "VAROMATIC+"
    subtitle.text = "Next-Generation Football Analysis System"
    
    # Overview slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Overview'
    
    tf = body_shape.text_frame
    tf.text = 'Key Features:'
    
    p = tf.add_paragraph()
    p.text = '• Real-time player detection'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = '• Automated offside detection'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = '• Team color analysis'
    p.level = 1
    
    # Technical slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Technical Stack'
    
    tf = body_shape.text_frame
    tf.text = 'Technologies:'
    
    p = tf.add_paragraph()
    p.text = '• YOLOv8 Neural Network'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = '• OpenCV'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = '• PyQt6'
    p.level = 1
    
    # Save the presentation
    prs.save('VAROMATIC+_Simple.pptx')

if __name__ == '__main__':
    create_presentation()
    print("Presentation created successfully!") 