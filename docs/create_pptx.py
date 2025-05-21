from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "VAROMATIC+"
    subtitle.text = "Next-Generation Football Analysis System"
    
    # Style the text
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)

def create_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Overview"
    
    tf = content.text_frame
    tf.text = "Key Features:"
    
    p = tf.add_paragraph()
    p.text = "• State-of-the-art football analysis system"
    p = tf.add_paragraph()
    p.text = "• Real-time player detection and tracking"
    p = tf.add_paragraph()
    p.text = "• Automated offside detection"
    p = tf.add_paragraph()
    p.text = "• Team color analysis"
    p = tf.add_paragraph()
    p.text = "• Professional-grade UI"

def create_technical_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Stack"
    
    tf = content.text_frame
    tf.text = "Core Technologies:"
    
    p = tf.add_paragraph()
    p.text = "• YOLOv8 Neural Network for player detection"
    p = tf.add_paragraph()
    p.text = "• OpenCV for video processing"
    p = tf.add_paragraph()
    p.text = "• PyQt6 for professional UI"
    p = tf.add_paragraph()
    p.text = "• Custom color analysis algorithms"
    p = tf.add_paragraph()
    p.text = "• Advanced offside detection system"

def create_performance_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Performance Metrics"
    
    tf = content.text_frame
    tf.text = "Real-time Analysis:"
    
    p = tf.add_paragraph()
    p.text = "• 30 FPS processing speed"
    p = tf.add_paragraph()
    p.text = "• 95% player detection accuracy"
    p = tf.add_paragraph()
    p.text = "• GPU-accelerated processing"
    p = tf.add_paragraph()
    p.text = "• Efficient memory management"
    p = tf.add_paragraph()
    p.text = "• Multi-threaded operation"

def create_future_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Developments"
    
    tf = content.text_frame
    tf.text = "Upcoming Features:"
    
    p = tf.add_paragraph()
    p.text = "• Enhanced player identification"
    p = tf.add_paragraph()
    p.text = "• Tactical analysis system"
    p = tf.add_paragraph()
    p.text = "• Cloud integration"
    p = tf.add_paragraph()
    p.text = "• Mobile application"
    p = tf.add_paragraph()
    p.text = "• API services"

def create_contact_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Contact Information"
    
    tf = content.text_frame
    tf.text = "VAROMATIC+ Professional"
    
    p = tf.add_paragraph()
    p.text = "Website: www.varomatic.pro"
    p = tf.add_paragraph()
    p.text = "Email: support@varomatic.pro"
    p = tf.add_paragraph()
    p.text = "Twitter: @varomaticpro"

def main():
    # Create presentation
    prs = Presentation()
    
    # Create slides
    create_title_slide(prs)
    create_overview_slide(prs)
    create_technical_slide(prs)
    create_performance_slide(prs)
    create_future_slide(prs)
    create_contact_slide(prs)
    
    # Save presentation
    prs.save('VAROMATIC+_Presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    main() 